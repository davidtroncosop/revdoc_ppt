import os
import tempfile
import streamlit as st
from dotenv import load_dotenv
import zipfile
import fitz  # PyMuPDF
import pandas as pd
from PIL import Image as PILImage
import io
import requests
import base64
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Cargar variables de entorno desde .env
load_dotenv()
api_key = os.getenv('OPENAI_API_KEY')

if not api_key:
    st.error("API Key no encontrada. Asegúrate de que el archivo .env esté correctamente configurado.")
    st.stop()

client = OpenAI(api_key=api_key)

# Funciones para analizar imágenes y documentos PDF
def encode_image(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

def analyze_image(base64_image, api_key):
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}"
    }

    payload = {
        "model": "gpt-4o-mini",
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": "Eres un asistente virtual especializado en la revisión de documentos escaneados. Analiza la imagen adjunta y extrae la siguiente información de manera breve y estructurada:\n\n- Tipo de documento (incluyendo documentos de identidad)\n- Nombres completos\n- Fechas relevantes\n- Institución emisora\n- Diagnóstico médico (si aplica)\n- Firmas y sellos presentes\n- Resumen de la carta** (si el documento es una carta o contiene una carta, proporciona un resumen conciso de su contenido)\n\nAdjunto una imagen para que la revises."
                        },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/jpeg;base64,{base64_image}"
                        }
                    }
                ]
            }
        ],
        "max_tokens": 300
    }

    response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload)
    
    if response.status_code != 200:
        st.error(f"Error: {response.status_code} - {response.text}")
        return None
    
    response_data = response.json()
    if 'error' in response_data:
        st.error(f"API Error: {response_data['error']['message']}")
    
    return response_data

def convert_pdf_page_to_image(pdf_path, page_num):
    pdf_document = fitz.open(pdf_path)
    page = pdf_document.load_page(page_num)
    pix = page.get_pixmap()
    img = PILImage.frombytes("RGB", [pix.width, pix.height], pix.samples)
    return img

def process_pdfs_in_zip(zip_path, output_dir, api_key):
    with zipfile.ZipFile(zip_path, 'r') as zip_ref:
        zip_ref.extractall(output_dir)

    results = []

    for root, dirs, files in os.walk(output_dir):
        for file in files:
            if file.endswith('.pdf'):
                pdf_path = os.path.join(root, file)
                folder_name = os.path.basename(os.path.dirname(pdf_path))
                st.write(f"Processing {pdf_path}...")

                pdf_document = fitz.open(pdf_path)
                pdf_analysis = []
                for page_num in range(len(pdf_document)):
                    img = convert_pdf_page_to_image(pdf_path, page_num)
                    temp_img_path = os.path.join(root, f"temp_image_page_{page_num}.png")
                    img.save(temp_img_path)

                    base64_image = encode_image(temp_img_path)
                    result = analyze_image(base64_image, api_key)
                    
                    if result and 'choices' in result and len(result['choices']) > 0:
                        analysis = result['choices'][0]['message']['content']
                        pdf_analysis.append(analysis)
                    else:
                        pdf_analysis.append("Analysis failed")

                results.append({
                    "Folder": folder_name,
                    "Analysis": " ".join(pdf_analysis)
                })


    # Convertir los resultados a un DataFrame
    data = pd.DataFrame(results)

    # Asegurarse de que 'Folder' sea de tipo string
    data['Folder'] = data['Folder'].astype(str)

    # Agrupar los datos por 'Folder' y concatenar la columna 'Analysis'
    result = data.groupby('Folder')['Analysis'].apply(lambda x: ' '.join(x)).reset_index()

    # Renombrar la columna para mayor claridad
    result.columns = ['Folder', 'Análisis_concatenado']

    return result

def merge_analysis_with_excel(excel_path, pdf_analysis_results):
    # Leer el archivo Excel
    data2 = pd.read_excel(excel_path)

    # Renombrar la columna 'RUT:' a 'Folder' si existe
    if 'RUT:' in data2.columns:
        data2.rename(columns={'RUT:': 'Folder'}, inplace=True)

    # Asegurarse de que 'Folder' sea de tipo string en ambos DataFrames
    data2['Folder'] = data2['Folder'].astype(str)
    pdf_analysis_results['Folder'] = pdf_analysis_results['Folder'].astype(str)

    # Realizar la unión de los dos DataFrames
    merged_data = pd.merge(data2, pdf_analysis_results, on='Folder', how='left')
    
    return merged_data


def generar_propuesta_resolucion(filas):
    resultados = []
    for _, fila in filas.iterrows():
        prompt = f"""
A continuación se presenta la información de un estudiante. Con base en esta información, por favor genera una "Propuesta Resolución" que indique si se aprueba o rechaza la solicitud de beca, y los detalles de la resolución. Usa las siguientes condiciones para tomar la decisión:

1. Si el PPE es menor a 0.5, rechaza la solicitud porque no cumple con el requisito mínimo.
2. Si la deuda vencida en el sistema es 0, rechaza la solicitud porque no hay deuda a cubrir.
3. Si los documentos han sido validados por una Trabajadora Social y el estudiante tiene una deuda vencida mayor que 0, aprueba la solicitud con los detalles correspondientes.
4. Si el estudiante ha recibido beneficios anteriormente, verifica si hay algún incumplimiento relacionado y decide en consecuencia.

Información del Estudiante:
PPE: {fila['PPE']}
Nombre completo: {fila['Nombre completo']}
RUT: {fila['Folder']}
Sede: {fila['Sede']}
Carrera: {fila['Carrera']}
Vigencia con cursos inscritos: {fila['Vigencia con cursos inscritos']}
Año y Semestre de ingreso: {fila['Año y Semestre de ingreso']}
Motivo solicitud: {fila['Motivo solicitud.']}
¿Ha recibido beneficios anteriormente? ¿Cuál?: {fila['¿Ha recibido beneficios anteriormente? ¿Cuál?']}
Última fecha en que se entregó el Beneficio: {fila['Última fecha en que se entregó el Beneficio']}
Deuda vencida en sistema: {fila['Deuda vencida en sistema']}
Análisis_concatenado: {fila['Análisis_concatenado']}

Por favor, genera una respuesta en formato de tabla con las siguientes columnas, incluyendo el encabezado, usando guiones (-) como delimitadores:
Propuesta Resolución-RESOLUCIÓN-MONTO DE LA BECA-MOTIVO DEL CASO-DOCUMENTOS

**La columna "RESOLUCIÓN" debe contener la decisión tomada sobre la solicitud (aprobada o rechazada) y la justificación basada en los criterios establecidos.**

**La columna "MOTIVO DEL CASO" debe contener un resumen muy breve de la situación presentada por el estudiante en su carta de solicitud, incluyendo detalles como problemas económicos, personales o académicos mencionados. Esta columna no debe repetir el contenido de la "RESOLUCIÓN", sino que debe reflejar la situación específica del estudiante tal como se describe en la carta.**

**En la columna "MOTIVO DEL CASO",  proporciona una descripción detallada en el siguiente formato: "Se informa lo siguiente: >[Detalle1]. >[Detalle2]. >[Detalle3].". Asegúrate de incluir todos los detalles relevantes separados por punto y coma (;), sin espacios adicionales.**

Ejemplo:
Propuesta Resolución-RESOLUCIÓN-MONTO DE LA BECA-MOTIVO DEL CASO-DOCUMENTOS
Aprobada-La solicitud de beca se aprueba...-Monto a determinar según normativa-Los Documentos informan lo siguiente >Familia extensa. >Un integrante genera ingresos formales (pensión). >El estudiante no encuentra empleo y recibe ayuda económica de su tío y abuela. >Postulación FUAS: octubre de 2022. >No presenta resultados MINEDUC. >Se sugiere aprobar la solicitud.-Carta de solicitud de beca; Registro Social de Hogares; Certificado de remuneraciones; Finiquitos; Certificado de cotizaciones; Licencias medicas; Comprobante de gastos mensuales; Certificado de desempleo
Rechazada-La solicitud de beca se rechaza porque...-Documentos validados por la Trabajadora Social, quien informa lo siguiente: >Familia extensa. >Un integrante genera ingresos formales (pensión). >El estudiante no encuentra empleo y recibe ayuda económica de su tío y abuela. >Postulación FUAS: octubre de 2022. >No presenta resultados MINEDUC. >Se sugiere aprobar la solicitud.-Carta de solicitud de beca; Registro Social de Hogares; Certificado de remuneraciones; Finiquitos; Certificado de cotizaciones; Licencias medicas; Comprobante de gastos mensuales; Certificado de desempleo

**Asegúrate de que cada valor esté correctamente delimitado por guiones y de que no haya espacios adicionales antes o después de los guiones. Incluye solo las 5 columnas especificadas y usa punto y coma para separar múltiples documentos en la columna DOCUMENTOS.** 
""" 


        
        completion = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "Eres un asistente social."},
                {"role": "user", "content": prompt}
            ]
        )

        respuesta_fila = completion.choices[0].message.content.strip()
        print(f"Raw GPT response: {respuesta_fila}")

        # Eliminar el encabezado de la respuesta
        if 'Propuesta Resolución' in respuesta_fila:
            respuesta_fila = respuesta_fila.split('\n')[1]

        # Dividir la fila en columnas usando el guion "-"
        columnas = respuesta_fila.split('-')

        # Si hay más de 5 columnas, combinar las columnas extra en la última columna
        if len(columnas) > 5:
            columnas[4] = '-'.join(columnas[4:])  # Combina columnas extra en la columna de DOCUMENTOS
            columnas = columnas[:5]  # Mantén solo las primeras 5 columnas

        # Si hay menos columnas de las esperadas, agregar columnas vacías
        while len(columnas) < 5:
            columnas.append('')

        # Convertir las columnas a string
        columnas = [str(valor).strip() for valor in columnas]

        # Agregar las columnas a la lista de resultados
        resultados.append(columnas)

    return resultados

def add_textbox(slide, left, top, width, height, text, font_size=Pt(14), bold=False, font_color=RGBColor(0, 0, 0), alignment=PP_ALIGN.LEFT):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = font_color
    p.alignment = alignment

def create_header_background(slide, left, top, width, height):
    background = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 192, 0)  # Light blue
    background.line.color.rgb = RGBColor(142, 180, 227)  # Sky blue border

def create_card(slide, left, top, width, height, title, subtitles, contents, image_path=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
    card.line.color.rgb = RGBColor(200, 200, 200)  # Light gray border

    # Add title
    title_box = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.25), width - Inches(0.5), Inches(0.5))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(18)
    title_box.text_frame.paragraphs[0].font.bold = True

    # Add subtitles and contents
    content_top = top + Inches(0.75)
    for subtitle, content in zip(subtitles, contents):
        subtitle_box = slide.shapes.add_textbox(left + Inches(0.25), content_top, width - Inches(0.5), Inches(0.3))
        subtitle_box.text_frame.text = subtitle
        subtitle_box.text_frame.paragraphs[0].font.size = Pt(14)
        subtitle_box.text_frame.paragraphs[0].font.bold = True
        content_top += Inches(0.3)

        content_box = slide.shapes.add_textbox(left + Inches(0.25), content_top, width - Inches(0.5), Inches(1))
        content_box.text_frame.word_wrap = True
        if isinstance(content, str):
            content_box.text_frame.text = content
        else:
            content_box.text_frame.text = str(content)
        for paragraph in content_box.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
        content_top += Inches(1)

def create_button(slide, left, top, width, height, text, color):
    button = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    button.fill.solid()
    button.fill.fore_color.rgb = color
    button.line.color.rgb = color
    button.text_frame.text = text
    button.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text


def create_slide_from_row(prs, row):
    # Create slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add header background
    create_header_background(slide, Inches(0.25), Inches(0.2), Inches(12.9), Inches(1.6))

    # Add header information
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(4), Inches(0.5),
                f"{row.get('Nombre completo', 'N/A')}\n{row.get('Folder', 'N/A')}",
                font_size=Pt(14), bold=True)

    add_textbox(slide, Inches(5), Inches(0.2), Inches(4), Inches(0.5),
                f"{row.get('Carrera', 'N/A')}\n{row.get('Sede', 'N/A')}",
                font_size=Pt(14), bold=True)

    add_textbox(slide, Inches(9.5), Inches(0.2), Inches(3.5), Inches(0.5),
                f"MATRÍCULA 2024-1\nCON CURSOS: {row.get('Vigencia con cursos inscritos', 'N/A')}",
                font_size=Pt(14), bold=True, alignment=PP_ALIGN.RIGHT)

    # Add timeline
    add_textbox(slide, Inches(0.7), Inches(1), Inches(2), Inches(0.3),
                f"Ingresa: {row.get('Año y Semestre de ingreso', 'N/A')}",
                font_size=Pt(10), font_color=RGBColor(255, 255, 255))
    

    add_textbox(slide, Inches(10.5), Inches(1), Inches(2.5), Inches(0.3),
                f"Envía Solicitud: {row.get('Hora de inicio', 'N/A')}",
                font_size=Pt(10), font_color=RGBColor(255, 255, 255), alignment=PP_ALIGN.RIGHT)

    # Create three cards
    card_width = Inches(4)
    card_height = Inches(5.3)
    spacing = Inches(0.5)

    # First card content
    subtitles_list1 = ["MOTIVO", "MOTIVO DEL CASO","DOCUMENTOS"]
    contents_list1 = [
        row.get('Motivo solicitud.', 'N/A'),
        row.get('MOTIVO DEL CASO', 'No hay información disponible'),
        row.get('DOCUMENTOS', 'No hay información disponible')
    ]

    card_left = Inches(0.1) + 0 * (card_width + spacing)
    card_top = Inches(2)
    create_card(slide, card_left, card_top, card_width, card_height,
                "SOLICITA", subtitles_list1, contents_list1)

    # Extract FUAS information safely
    analisis = str(row.get('Análisis_concatenado', ''))
    fuas_info = 'No especificado'
    if 'Postulación FUAS:' in analisis:
        fuas_info = analisis.split('Postulación FUAS:')[1].split('.')[0].strip()

    subtitles_list2 = ["ANTECEDENTES ECONÓMICOS", "ANTECEDENTES ACADÉMICOS"]
    contents_list2 = [
        f"Beneficio: {row.get('¿Ha recibido beneficios anteriormente? ¿Cuál?', 'N/A')}\n"
        f"Deuda: ${'{:,}'.format(row.get('Deuda vencida en sistema', 0.0))}\n"
        f"Postulación a FUAS: {fuas_info}\n"
        f"Arancel: ${'{:,}'.format(row.get('Monto cuota de Arancel', 0.0))}\n"
        f"Matrícula: ${'{:,}'.format(row.get('Monto valor de matrícula', 0.0))}",
        f"Avance Curricular: {row.get('Avance curricular (%)', 'N/A')}\n"
        f"PPS: {row.get('PPS', 'N/A')}\n"
        f"RSH: {row.get('Registro Social de Hogares (RSH) o Nivel Socioeconómico (NSE)', 'N/A')}\n"
        f"Promedio Ponderado Evaluación: {'{:.2f}'.format(row.get('PPE', 0.0))}"
    ]
    
    # Extract resolution information safely
    resolucion = row.get('RESOLUCIÓN', 'No hay información disponible')
    monto_beca = f"${'{:,}'.format(row.get('Plan de Retención', 0.0))}"
     
    subtitles_list3 = ["RESOLUCIÓN", "MONTO DE LA BECA"]
    contents_list3 = [resolucion, monto_beca]

    button_left = card_left + Inches(0.25)
    button_top = card_top + card_height - Inches(0.7)
    button_width = card_width - Inches(0.5)
    button_height = Inches(0.5)
    create_button(slide, button_left, button_top, button_width, button_height,
                  "Solicitud", RGBColor(13, 34, 60))  # Indigo color

    card_left = Inches(0.1) + 1 * (card_width + spacing)
    create_card(slide, card_left, card_top, card_width, card_height,
                "ANTECEDENTES", subtitles_list2, contents_list2)

    button_left = card_left + Inches(0.25)
    create_button(slide, button_left, button_top, button_width, button_height,
                  "Revisión", RGBColor(13, 34, 60))  # Teal color

    card_left = Inches(0.1) + 2 * (card_width + spacing)
    create_card(slide, card_left, card_top, card_width, card_height,
                "RESOLUCIÓN", subtitles_list3, contents_list3)

    button_left = card_left + Inches(0.25)
    create_button(slide, button_left, button_top, button_width, button_height,
                  "Resolución", RGBColor(13, 34, 60))   # Blue color

def create_presentation_from_dataframe(df, output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for _, row in df.iterrows():
        create_slide_from_row(prs, row)

    prs.save(output_path)

# Interfaz de usuario Streamlit
st.title("Análisis de documentos PDF y generación de propuestas")

uploaded_zip = st.file_uploader("Sube un archivo ZIP con PDFs", type=["zip"])
uploaded_excel = st.file_uploader("Sube un archivo Excel con datos de estudiantes", type=["xlsx"])

if uploaded_zip and uploaded_excel:
    with tempfile.TemporaryDirectory() as temp_dir:
        zip_path = os.path.join(temp_dir, "uploaded.zip")
        with open(zip_path, "wb") as zip_file:
            zip_file.write(uploaded_zip.getbuffer())

        excel_path = os.path.join(temp_dir, "uploaded.xlsx")
        with open(excel_path, "wb") as excel_file:
            excel_file.write(uploaded_excel.getbuffer())

        pdf_analysis_results = process_pdfs_in_zip(zip_path, temp_dir, api_key)
        
        # Fusionar el análisis de PDF con el DataFrame del Excel
        df = merge_analysis_with_excel(excel_path, pdf_analysis_results)
        
        # Guardar el DataFrame actualizado con el análisis
        df.to_excel(os.path.join(temp_dir, "excel_con_analisis.xlsx"), index=False)

        st.success("Análisis de PDFs completado y fusionado con el Excel original.")



        propuestas = generar_propuesta_resolucion(df)
        df_propuestas = pd.DataFrame(propuestas, columns=["Propuesta Resolución", "RESOLUCIÓN", "MONTO DE LA BECA", "MOTIVO DEL CASO", "DOCUMENTOS"])
        df_final = pd.concat([df, df_propuestas], axis=1)
        df_final.to_excel(os.path.join(temp_dir, "resultado_final.xlsx"), index=False)

        st.success("Propuesta de resolución generada y guardada en resultado_final.xlsx")

        pptx_path = os.path.join(temp_dir, "presentacion_estudiantes.pptx")
        create_presentation_from_dataframe(df_final, pptx_path)

        st.success("Presentación PowerPoint generada y guardada en presentacion_estudiantes.pptx")

        # Ofrecer los archivos para descargar
        with open(os.path.join(temp_dir, "excel_con_analisis.xlsx"), "rb") as file:
            st.download_button(
                label="Descargar Excel con análisis de PDFs",
                data=file,
                file_name="excel_con_analisis.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )

        with open(os.path.join(temp_dir, "resultado_final.xlsx"), "rb") as file:
            st.download_button(
                label="Descargar resultado final",
                data=file,
                file_name="resultado_final.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )

        with open(pptx_path, "rb") as file:
            st.download_button(
                label="Descargar presentación PowerPoint",
                data=file,
                file_name="presentacion_estudiantes.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )